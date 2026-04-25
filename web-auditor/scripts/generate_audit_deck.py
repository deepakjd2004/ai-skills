import argparse
import json
import datetime
from pathlib import Path
from urllib.parse import urlparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# --- Configuration ---
JSON_FILE = "../output/performance_audit_www_example_com_20260424_153606.json"
OUTPUT_DIR = "../output/"

# Severity color mapping
COLORS = {
    "CRITICAL": RGBColor(255, 68, 68),   # Red
    "HIGH": RGBColor(255, 107, 53),      # Orange
    "MEDIUM": RGBColor(255, 179, 71),    # Amber
    "DEFAULT": RGBColor(0, 153, 204)     # Cyan
}


def parse_args():
    parser = argparse.ArgumentParser(
        description="Generate a performance audit PowerPoint from an audit JSON file."
    )
    parser.add_argument(
        "--json-file",
        default=JSON_FILE,
        help="Path to the audit JSON file.",
    )
    parser.add_argument(
        "--template-file",
        default=None,
        help=(
            "Optional path to a PPTX template. "
            "If omitted, a standard presentation is generated without a template."
        ),
    )
    parser.add_argument(
        "--output-dir",
        default=OUTPUT_DIR,
        help="Directory where the generated PPTX file will be written.",
    )
    return parser.parse_args()

def get_rating(metric, value):
    """Returns the CWV rating based on standard Google thresholds."""
    if value is None:
        return "N/A"
    
    thresholds = {
        "LCP": (2500, 4000),
        "INP": (200, 500),
        "CLS": (0.10, 0.25),
        "FCP": (1800, 3000),
        "TTFB": (800, 1800)
    }
    
    good, poor = thresholds.get(metric, (0, 0))
    val = float(value)
    
    if val <= good:
        return "Good"
    elif val <= poor:
        return "Needs Improvement"
    else:
        return "Poor"


def parse_timestamp(value):
    """Parse ISO timestamp values, including trailing Z."""
    if value.endswith("Z"):
        value = value[:-1] + "+00:00"
    return datetime.datetime.fromisoformat(value)


def extract_cwv_metrics(crux):
    return [
        ("LCP (Largest Contentful Paint)", "LCP", crux.get("largest_contentful_paint", {}).get("p75")),
        ("INP (Interaction to Next Paint)", "INP", crux.get("interaction_to_next_paint", {}).get("p75")),
        ("CLS (Cumulative Layout Shift)", "CLS", crux.get("cumulative_layout_shift", {}).get("p75")),
        ("FCP (First Contentful Paint)", "FCP", crux.get("first_contentful_paint", {}).get("p75")),
        ("TTFB (Time to First Byte)", "TTFB", crux.get("experimental_time_to_first_byte", {}).get("p75")),
    ]


def group_recommendations(recommendations):
    grouped = {}
    for rec in recommendations:
        category = rec.get("category", "General")
        grouped.setdefault(category, []).append(rec)
    return grouped

def move_slide(prs, old_index, new_index):
    """Moves a slide in the presentation from old_index to new_index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def replace_paragraph_text(para, replacements):
    """Replace text patterns in a paragraph, handling cross-run splits.

    python-pptx can split a single visible string (e.g. '03.2026') across
    multiple runs with different formatting, so per-run checks miss matches
    that span a run boundary.  This helper joins all run text, applies every
    replacement, then writes the result back into the first run and blanks
    the rest.
    """
    if not para.runs:
        return
    full_text = "".join(run.text for run in para.runs)
    new_text = full_text
    for old, new in replacements:
        new_text = new_text.replace(old, new)
    if new_text != full_text:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""


def build_cwv_table(slide, metrics):
    table_shape = slide.shapes.add_table(
        len(metrics) + 1, 4, Inches(0.45), Inches(1.35), Inches(9.1), Inches(4.75)
    )
    table = table_shape.table

    table.columns[0].width = Inches(3.8)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(2.6)

    header_bg = COLORS["DEFAULT"]
    header_labels = ["Metric", "Code", "Value", "Status"]
    for col_idx, label in enumerate(header_labels):
        cell = table.cell(0, col_idx)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, (label, code, val) in enumerate(metrics, start=1):
        rating = get_rating(code, val)
        unit = "" if code == "CLS" else " ms"
        value_text = "N/A" if val is None else f"{val}{unit}"

        if rating == "Good":
            status_bg = RGBColor(46, 204, 113)
        elif rating == "Needs Improvement":
            status_bg = RGBColor(255, 193, 7)
        else:
            status_bg = RGBColor(255, 82, 82)

        row_bg = RGBColor(237, 248, 255) if row_idx % 2 else RGBColor(225, 242, 252)

        metric_cell = table.cell(row_idx, 0)
        metric_cell.text = label
        metric_cell.fill.solid()
        metric_cell.fill.fore_color.rgb = row_bg
        p_metric = metric_cell.text_frame.paragraphs[0]
        p_metric.font.bold = True
        p_metric.font.size = Pt(12)
        p_metric.font.color.rgb = RGBColor(28, 48, 64)

        code_cell = table.cell(row_idx, 1)
        code_cell.text = code
        code_cell.fill.solid()
        code_cell.fill.fore_color.rgb = row_bg
        p_code = code_cell.text_frame.paragraphs[0]
        p_code.font.bold = True
        p_code.font.size = Pt(12)
        p_code.font.color.rgb = COLORS["DEFAULT"]

        value_cell = table.cell(row_idx, 2)
        value_cell.text = value_text
        value_cell.fill.solid()
        value_cell.fill.fore_color.rgb = row_bg
        p_value = value_cell.text_frame.paragraphs[0]
        p_value.font.size = Pt(12)
        p_value.font.color.rgb = RGBColor(20, 20, 20)

        status_cell = table.cell(row_idx, 3)
        status_cell.text = rating
        status_cell.fill.solid()
        status_cell.fill.fore_color.rgb = status_bg
        p_status = status_cell.text_frame.paragraphs[0]
        p_status.font.bold = True
        p_status.font.size = Pt(12)
        p_status.font.color.rgb = RGBColor(255, 255, 255)


def add_recommendations_to_text_frame(text_frame, items, light_mode=False):
    text_frame.clear()
    for item in items:
        severity = item.get("severity", "DEFAULT")
        issue = item.get("issue", "Unknown Issue")
        recommendation = item.get("recommendation", "")

        p_issue = text_frame.add_paragraph()
        p_issue.level = 0
        p_issue.space_before = Pt(10)

        run_sev = p_issue.add_run()
        run_sev.text = f"{severity}: "
        run_sev.font.bold = True
        run_sev.font.color.rgb = COLORS.get(severity, COLORS["DEFAULT"])

        run_issue = p_issue.add_run()
        run_issue.text = issue

        p_rec = text_frame.add_paragraph()
        p_rec.level = 1
        p_rec.space_after = Pt(10)
        run_rec = p_rec.add_run()
        run_rec.text = recommendation
        if light_mode:
            run_rec.font.color.rgb = RGBColor(40, 40, 40)
        else:
            run_rec.font.color.rgb = RGBColor(180, 255, 200)


def get_content_shape(slide):
    if len(slide.placeholders) > 1:
        return slide.placeholders[1]
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            return shape
    return None


def build_template_presentation(prs, audit, month_str, year_str, current_date_str, metrics, grouped_recs):
    url = audit["url"]

    # Cover slide updates
    if len(prs.slides) > 0:
        slide_cover = prs.slides[0]
        for shape in slide_cover.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replacements = [
                    ("https://www.example.com/", url),
                    ("https://www.example.com", url),
                    ("example.com", url),
                    ("03.", f"{month_str}."),
                    ("2026", year_str),
                    ("<DATE>", current_date_str),
                ]
                replace_paragraph_text(paragraph, replacements)
                for run in paragraph.runs:
                    if url in run.text:
                        if len(url) > 60:
                            run.font.size = Pt(14)
                        elif len(url) > 40:
                            run.font.size = Pt(18)
                        elif len(url) > 25:
                            run.font.size = Pt(24)

    # CWV table goes to slide 5 (index 4) when available; otherwise append.
    if len(prs.slides) > 4:
        slide_cwv = prs.slides[4]
    else:
        slide_cwv = prs.slides.add_slide(prs.slide_layouts[5])
        if slide_cwv.shapes.title:
            slide_cwv.shapes.title.text = "Core Web Vitals"
    build_cwv_table(slide_cwv, metrics)

    # Observation slides from template layout when available, else append using title+content.
    if len(prs.slides) > 8:
        obs_layout = prs.slides[8].slide_layout
        insert_index = min(9, len(prs.slides))

        for category, items in grouped_recs.items():
            new_slide = prs.slides.add_slide(obs_layout)
            title_shape = new_slide.shapes.title
            if title_shape:
                title_shape.text = f"Observations: {category}"
                para = title_shape.text_frame.paragraphs[0]
                if para.runs:
                    para.runs[0].font.color.rgb = COLORS["DEFAULT"]

                pPr = para._p.get_or_add_pPr()
                for tag in (
                    qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum"), qn("a:buBlip"),
                    qn("a:buFont"), qn("a:buSzPts"), qn("a:buClr")
                ):
                    for el in pPr.findall(tag):
                        pPr.remove(el)
                pPr.append(etree.SubElement(pPr, qn("a:buNone")))

                tx_body = title_shape.text_frame._txBody
                lst_style = tx_body.find(qn("a:lstStyle"))
                if lst_style is None:
                    lst_style = etree.SubElement(tx_body, qn("a:lstStyle"))
                lvl1 = lst_style.find(qn("a:lvl1pPr"))
                if lvl1 is None:
                    lvl1 = etree.SubElement(lst_style, qn("a:lvl1pPr"))
                for tag in (
                    qn("a:buNone"), qn("a:buChar"), qn("a:buAutoNum"), qn("a:buBlip"),
                    qn("a:buFont"), qn("a:buSzPts"), qn("a:buClr")
                ):
                    for el in lvl1.findall(tag):
                        lvl1.remove(el)
                lvl1.append(etree.SubElement(lvl1, qn("a:buNone")))

            body_shape = get_content_shape(new_slide)
            if body_shape is not None:
                add_recommendations_to_text_frame(body_shape.text_frame, items, light_mode=False)

            current_index = len(prs.slides) - 1
            move_slide(prs, current_index, insert_index)
            insert_index += 1
    else:
        for category, items in grouped_recs.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Observations: {category}"
            body_shape = get_content_shape(slide)
            if body_shape is not None:
                add_recommendations_to_text_frame(body_shape.text_frame, items, light_mode=True)


def build_standard_presentation(prs, audit, current_date_str, metrics, grouped_recs):
    url = audit["url"]
    technical = audit.get("technical_checks", {})

    # Slide 1: Cover
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = "Web Performance Review"
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = f"Target: {url}\nAudit Date: {current_date_str}"

    # Slide 2: Technical summary
    summary = prs.slides.add_slide(prs.slide_layouts[1])
    summary.shapes.title.text = "Technical Summary"
    body = get_content_shape(summary)
    if body is not None:
        tf = body.text_frame
        tf.clear()
        points = [
            f"Primary HTTP version: {technical.get('http_version', 'Unknown')}",
            f"IPv6 support: {technical.get('ipv6_support', 'Unknown')}",
            f"Cache policy: {technical.get('cache_analysis', {}).get('summary', 'Unknown')}",
        ]
        for idx, point in enumerate(points):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = point

    # Slide 3: CWV table
    cwv = prs.slides.add_slide(prs.slide_layouts[5])
    if cwv.shapes.title:
        cwv.shapes.title.text = "Core Web Vitals"
    build_cwv_table(cwv, metrics)

    # Slides 4+: grouped recommendations
    for category, items in grouped_recs.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Observations: {category}"
        body_shape = get_content_shape(slide)
        if body_shape is not None:
            add_recommendations_to_text_frame(body_shape.text_frame, items, light_mode=True)

def main():
    args = parse_args()

    json_path = Path(args.json_file).expanduser().resolve()
    output_dir = Path(args.output_dir).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    with json_path.open("r", encoding="utf-8") as f:
        audit = json.load(f)

    url = audit.get("url", "unknown-url")
    domain = urlparse(url).netloc or "unknown_domain"

    dt = parse_timestamp(audit["timestamp"])
    current_date_str = datetime.datetime.now().strftime("%d-%m-%Y")
    month_str = dt.strftime("%m")
    year_str = dt.strftime("%Y")

    crux = audit.get("crux_data", {})
    recommendations = audit.get("recommendations", [])
    metrics = extract_cwv_metrics(crux)
    grouped_recs = group_recommendations(recommendations)

    prs = None
    using_template = False
    template_path = Path(args.template_file).expanduser().resolve() if args.template_file else None

    if template_path is not None and template_path.exists():
        try:
            prs = Presentation(str(template_path))
            build_template_presentation(
                prs=prs,
                audit=audit,
                month_str=month_str,
                year_str=year_str,
                current_date_str=current_date_str,
                metrics=metrics,
                grouped_recs=grouped_recs,
            )
            using_template = True
        except Exception as exc:
            print(f"Template-based generation failed, falling back to standard PPT: {exc}")

    if prs is None:
        prs = Presentation()
        build_standard_presentation(
            prs=prs,
            audit=audit,
            current_date_str=current_date_str,
            metrics=metrics,
            grouped_recs=grouped_recs,
        )

    run_timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"Performance_Review_{domain.replace('.', '_')}_{year_str}{month_str}_{run_timestamp}.pptx"
    output_path = output_dir / output_filename
    prs.save(str(output_path))

    mode_text = "template" if using_template else "standard"
    print(f"Successfully generated {mode_text} presentation: {output_path}")

if __name__ == "__main__":
    main()
